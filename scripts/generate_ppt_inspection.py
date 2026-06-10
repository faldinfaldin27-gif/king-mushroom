"""Generate the Exercise 8 Task (2) PPT: AI Inspection Robot (engineering design report).

    python scripts/generate_ppt_inspection.py

Topic: an autonomous AI inspection robot for a king oyster mushroom (杏鮑菇) grow house,
reusing the YOLO26 vision system from Exercise 9. Follows the 4-stage engineering design
process (clarification -> conceptual -> embodiment -> detailed) + capstone (architecture,
BOM, economics, operation). Produces AI_Inspection_Robot_Ex8.pptx in the project root.
"""
import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "AI_Inspection_Robot_Ex8.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
LBLUE = RGBColor(0xD6, 0xE4, 0xFF)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
LGREEN = RGBColor(0xD7, 0xF0, 0xDB)
ORANGE = RGBColor(0xE8, 0x71, 0x0A)
LORANGE = RGBColor(0xFC, 0xE8, 0xD2)
RED = RGBColor(0xD2, 0x32, 0x2D)
GRAY = RGBColor(0x44, 0x44, 0x44)
LGRAY = RGBColor(0xEC, 0xEC, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def set_run(p, text, size, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return r


def bar(slide, color=BLUE, h=Inches(0.18)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def header(slide, title, tag=None):
    bar(slide, BLUE)
    tf = add_box(slide, Inches(0.5), Inches(0.28), Inches(10.6), Inches(0.9))
    set_run(tf.paragraphs[0], title, 28, NAVY, True)
    if tag:
        tt = add_box(slide, Inches(10.7), Inches(0.32), Inches(2.4), Inches(0.6))
        set_run(tt.paragraphs[0], tag, 13, WHITE, True, PP_ALIGN.CENTER)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.7),
                                    Inches(0.34), Inches(2.35), Inches(0.5))
        sh.fill.solid(); sh.fill.fore_color.rgb = ORANGE; sh.line.fill.background()
        sh.text_frame.word_wrap = True
        p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = tag; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = WHITE


def new(title, tag=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title, tag)
    return s


def box(slide, l, t, w, h, text, fill=LBLUE, line=BLUE, fsz=12, fcolor=NAVY,
        bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    tf = sh.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; r.font.size = Pt(fsz)
        r.font.bold = bold; r.font.color.rgb = fcolor
    return sh


def arrow(slide, l, t, w, h, down=False):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = GRAY; sh.line.fill.background()
    return sh


def bullets(slide, items, l=Inches(0.7), t=Inches(1.5), w=Inches(12), h=Inches(5.6)):
    tf = add_box(slide, l, t, w, h)
    for i, item in enumerate(items):
        txt, lvl, *style = item
        color = style[0] if style else GRAY
        bold = style[1] if len(style) > 1 else (lvl == 0)
        size = 22 - lvl * 3
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = "" if lvl == 0 else ("   • " if lvl == 1 else "        – ")
        set_run(p, prefix + txt, size, color, bold)
        p.space_after = Pt(6)
    return tf


def table(slide, headers, rows, l=Inches(0.7), t=Inches(1.6), w=Inches(11.9),
          rh=0.42, fsz=13, hsz=14, widths=None, highlight_last=False):
    nrow, ncol = len(rows) + 1, len(headers)
    tb = slide.shapes.add_table(nrow, ncol, l, t, w, Inches(rh * nrow)).table
    if widths:
        for j, ww in enumerate(widths):
            tb.columns[j].width = Inches(ww)
    for j, hh in enumerate(headers):
        c = tb.cell(0, j); c.text = hh
        r = c.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(hsz); r.font.bold = True; r.font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        last = highlight_last and (i == len(rows))
        for j, val in enumerate(row):
            c = tb.cell(i, j); c.text = str(val) if str(val) else " "
            r = c.text_frame.paragraphs[0].runs[0]
            r.font.size = Pt(fsz)
            if last:
                r.font.bold = True; r.font.color.rgb = WHITE
                c.fill.solid(); c.fill.fore_color.rgb = BLUE
            elif i % 2 == 0:
                c.fill.solid(); c.fill.fore_color.rgb = LGRAY
    return tb


def caption(slide, text):
    tf = add_box(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
    set_run(tf.paragraphs[0], text, 12, GRAY, align=PP_ALIGN.CENTER)


detect_img = None
for g in (sorted(glob.glob(str(ROOT / "outputs" / "detect" / "IMG2023*.jpg"))) +
          sorted(glob.glob(str(ROOT / "outputs" / "detect" / "*.jpg")))):
    detect_img = g; break

# ====================== SLIDES ======================

# 1. Title
s = prs.slides.add_slide(BLANK)
bar(s, BLUE, Inches(0.25))
b2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.25), SW, Inches(0.25))
b2.fill.solid(); b2.fill.fore_color.rgb = BLUE; b2.line.fill.background()
tf = add_box(s, Inches(1), Inches(1.9), Inches(11.3), Inches(1.6))
set_run(tf.paragraphs[0], "AI Inspection Robot", 46, NAVY, True, PP_ALIGN.CENTER)
p = tf.add_paragraph(); set_run(p, "Autonomous Patrol for a King Oyster Mushroom Farm", 24, RED, True, PP_ALIGN.CENTER)
tf2 = add_box(s, Inches(1), Inches(3.9), Inches(11.3), Inches(2.2))
set_run(tf2.paragraphs[0], "AIoT System Innovation Design — Exercise 8, Task (2)", 20, GRAY, False, PP_ALIGN.CENTER)
p = tf2.add_paragraph(); set_run(p, "Engineering Design Process: Clarification -> Conceptual -> Embodiment -> Detailed", 16, BLUE, False, PP_ALIGN.CENTER)
p = tf2.add_paragraph(); set_run(p, "Vision payload powered by the YOLO26 + SAM + LLM/VLM system from Exercise 9", 14, GRAY, False, PP_ALIGN.CENTER)

# 2. Design process overview
s = new("Engineering Design Process")
labels = [("1. Clarification\nof the Task", LORANGE, ORANGE),
          ("2. Conceptual\nDesign", LBLUE, BLUE),
          ("3. Embodiment\nDesign", LGREEN, GREEN),
          ("4. Detailed\nDesign", LORANGE, ORANGE)]
x = Inches(0.7)
for i, (lab, fill, line) in enumerate(labels):
    box(s, x, Inches(2.6), Inches(2.6), Inches(1.6), lab, fill, line, 16)
    x += Inches(2.7)
    if i < 3:
        arrow(s, x - Inches(0.12), Inches(3.15), Inches(0.35), Inches(0.5))
bullets(s, [
    ("Goal: turn a real farm need into a buildable, costed, deployable inspection robot.", 0, NAVY, True),
    ("Each stage narrows the solution space and adds engineering detail.", 0, GRAY, False),
    ("Capstone deliverables: System Architecture, Bill of Materials, Economic Evaluation, Final Operation.", 0, GRAY, False),
], t=Inches(4.7), h=Inches(2))

# 3. Stage 1 - background / problem
s = new("Problem Statement & Background", "Stage 1")
bullets(s, [
    ("Context: a king oyster mushroom (杏鮑菇) grow house holds hundreds of substrate bags (太空包) on multi-level racks under controlled climate.", 0, NAVY, True),
    ("Current pain points (manual inspection):", 0, NAVY, True),
    ("Counting mushrooms & estimating yield is slow, tiring, and inconsistent.", 1),
    ("Contamination / abnormal growth is found late, spreading loss between bags.", 1),
    ("Temperature, humidity and CO2 vary across racks but are spot-checked only.", 1),
    ("Harvest timing relies on human judgement of mushroom size.", 1),
    ("Opportunity:", 0, NAVY, True),
    ("An autonomous robot can patrol 24/7, run AI vision (YOLO26) on every rack, log climate, and push alerts to a cloud dashboard.", 1, GREEN, True),
])

# 4. Stage 1 - requirements list
s = new("Requirements List (Demands / Wishes)", "Stage 1")
table(s, ["#", "Requirement", "Type", "Target / Spec"], [
    ["1", "Autonomous navigation along grow-house aisles", "D", "SLAM + line follow, no manual driving"],
    ["2", "Detect & count mushrooms per bag", "D", "YOLO26, count error <= 10%"],
    ["3", "Measure mushroom size for harvest timing", "D", "+/- 2 mm with calibration"],
    ["4", "Monitor climate (T / RH / CO2)", "D", "log every rack, +/-0.5C, +/-3% RH"],
    ["5", "Detect contamination / anomaly", "W", "VLM flag + photo to dashboard"],
    ["6", "Operate in humid 18-25C environment", "D", "IP54, condensation tolerant"],
    ["7", "Battery endurance for one full patrol", "D", ">= 4 h, auto-dock charging"],
    ["8", "Safety around workers", "D", "obstacle stop, e-stop, < 0.6 m/s"],
    ["9", "Cloud reporting & alerts", "D", "MQTT / Wi-Fi, real-time"],
    ["10", "Low cost & maintainable", "W", "BOM < USD 3,000"],
], rh=0.45, fsz=12.5, widths=[0.5, 5.6, 0.9, 4.9])
caption(s, "D = Demand (must have)   W = Wish (desirable)")

# 5. Stage 2 - function structure
s = new("Function Structure", "Stage 2")
box(s, Inches(0.5), Inches(1.7), Inches(2.3), Inches(1.0), "Energy:\nbattery power", LORANGE, ORANGE, 12)
box(s, Inches(0.5), Inches(3.0), Inches(2.3), Inches(1.0), "Material:\nmushroom bags\n(environment)", LGREEN, GREEN, 11)
box(s, Inches(0.5), Inches(4.3), Inches(2.3), Inches(1.0), "Signal:\npatrol command", LBLUE, BLUE, 12)
box(s, Inches(3.4), Inches(2.4), Inches(3.3), Inches(2.6),
    "MAIN FUNCTION\nAutonomously inspect\nmushroom racks &\nreport status", RGBColor(0xFF,0xF2,0xCC), ORANGE, 15, NAVY)
subs = ["Move &\nnavigate", "Sense\nenvironment", "Capture &\nanalyse vision", "Communicate\n& store data"]
y = Inches(1.7)
for slab in subs:
    box(s, Inches(7.4), y, Inches(2.6), Inches(1.05), slab, LBLUE, BLUE, 12)
    arrow(s, Inches(6.8), y + Inches(0.25), Inches(0.5), Inches(0.45))
    y += Inches(1.25)
arrow(s, Inches(2.85), Inches(3.4), Inches(0.45), Inches(0.5))
caption(s, "Overall function decomposed into 4 sub-functions, each solved in the morphological matrix")

# 6. Stage 2 - morphological matrix
s = new("Morphological Matrix (Solution Principles)", "Stage 2")
table(s, ["Sub-function", "Option A", "Option B", "Option C (chosen *)"], [
    ["Locomotion", "Wheeled diff-drive *", "Tracked", "Rail-guided"],
    ["Navigation", "Line-follow only", "LiDAR SLAM + line *", "Full visual SLAM"],
    ["Vision compute", "Cloud only", "Edge (Jetson) *", "MCU only"],
    ["Camera", "RGB webcam", "RGB-D depth cam *", "Stereo pair"],
    ["Climate sensing", "DHT22 only", "SCD30 (T/RH/CO2) *", "Industrial probe"],
    ["Power", "Lead-acid", "LiFePO4 + auto-dock *", "Swappable Li-ion"],
    ["Comms", "Wi-Fi *", "4G LTE", "LoRa"],
], rh=0.5, fsz=12.5, widths=[2.4, 2.9, 2.9, 3.7])
caption(s, "* = principle selected for the chosen concept (Concept 2)")

# 7. Stage 2 - concept evaluation
s = new("Concept Evaluation & Selection", "Stage 2")
table(s, ["Criterion (weight)", "Concept 1\nRail-guided", "Concept 2\nWheeled+SLAM *", "Concept 3\nDrone"], [
    ["Coverage / flexibility (25%)", "2", "5", "4"],
    ["Vision quality near bags (20%)", "4", "5", "2"],
    ["Cost & maintenance (20%)", "3", "4", "2"],
    ["Endurance (15%)", "5", "4", "1"],
    ["Safety around workers (10%)", "5", "4", "2"],
    ["Install complexity (10%)", "1", "4", "3"],
    ["Weighted score", "3.05", "4.45 *", "2.55"],
], rh=0.5, fsz=12.5, widths=[3.5, 2.8, 2.8, 2.8], highlight_last=True)
caption(s, "Score 1-5. Concept 2 (wheeled differential-drive + LiDAR/line navigation) wins -> carried forward")

# 8. Stage 3 - layout
s = new("Preliminary Layout", "Stage 3")
box(s, Inches(4.6), Inches(1.55), Inches(4.1), Inches(0.9), "Sensor mast: RGB-D camera +\nLED light + pan-tilt", LBLUE, BLUE, 12)
box(s, Inches(4.6), Inches(2.65), Inches(4.1), Inches(0.9), "Edge AI box (Jetson) +\nclimate sensor (SCD30)", LGREEN, GREEN, 12)
box(s, Inches(4.6), Inches(3.75), Inches(4.1), Inches(0.9), "Mobile base: motors, encoders,\nIMU, LiDAR, MCU", LORANGE, ORANGE, 12)
box(s, Inches(4.6), Inches(4.85), Inches(4.1), Inches(0.9), "Battery (LiFePO4) +\npower management", RGBColor(0xFF,0xF2,0xCC), ORANGE, 12)
for yy in (2.45, 3.55, 4.65):
    arrow(s, Inches(6.4), Inches(yy), Inches(0.5), Inches(0.22), down=True)
bullets(s, [
    ("Vertical stack keeps", 0, NAVY, True),
    ("centre of gravity low", 1),
    ("camera at rack height", 1),
    ("sensors away from", 1),
    ("motor heat & vibration", 1),
    ("easy module swap for", 1),
    ("maintenance", 1),
], l=Inches(0.6), t=Inches(1.7), w=Inches(3.6), h=Inches(5))
caption(s, "Modular vertical layout: base -> power -> compute -> sensing")

# 9. Stage 3 - subsystem design
s = new("Subsystem Design", "Stage 3")
table(s, ["Subsystem", "Design choice", "Key function"], [
    ["Locomotion", "2x DC gear motors + encoders, caster", "Drive aisles, odometry"],
    ["Navigation", "2D LiDAR SLAM + floor line, IMU fusion", "Localize & path-follow"],
    ["Vision", "RGB-D cam on pan-tilt + LED ring", "Image each rack level"],
    ["AI compute", "NVIDIA Jetson Orin Nano (edge)", "Run YOLO26 / SAM / VLM"],
    ["Climate", "SCD30 (T/RH/CO2) + spare port", "Per-rack environment log"],
    ["Power", "24V LiFePO4 20Ah + auto-dock", "4h+ patrol, self-charge"],
    ["Control", "STM32 MCU (motor/safety) <-> Jetson", "Real-time + high-level split"],
    ["Comms", "Wi-Fi (MQTT) to cloud", "Stream data & alerts"],
], rh=0.48, fsz=12.5, widths=[2.1, 5.4, 4.4])

# 10. Stage 4 - component specs
s = new("Detailed Design — Component Specs", "Stage 4")
table(s, ["Component", "Selected part", "Spec"], [
    ["Edge AI", "Jetson Orin Nano 8GB", "40 TOPS, runs YOLO26n @ ~30 FPS"],
    ["Camera", "Intel RealSense D435i", "RGB-D, 0.3-3m depth, IMU"],
    ["LiDAR", "RPLiDAR A2M12", "12m, 360deg, 8000 samples/s"],
    ["MCU", "STM32F4 + motor driver", "PWM, encoder, e-stop logic"],
    ["Climate", "Sensirion SCD30", "CO2 +/-30ppm, T +/-0.3C, RH +/-2%"],
    ["Motors", "12V DC gearmotor + Hall encoder", "~0.6 m/s top speed"],
    ["Battery", "24V 20Ah LiFePO4 + BMS", "480 Wh, ~4-5 h runtime"],
    ["IMU", "BNO055 9-DOF", "Fused orientation"],
], rh=0.48, fsz=12.5, widths=[2.0, 4.0, 5.9])

# 11. Stage 4 - electrical / wiring
s = new("Electrical / Wiring Diagram", "Stage 4")
box(s, Inches(0.5), Inches(3.0), Inches(2.2), Inches(1.1), "24V LiFePO4\nBattery + BMS", RGBColor(0xFF,0xF2,0xCC), ORANGE, 12)
box(s, Inches(3.0), Inches(3.0), Inches(2.2), Inches(1.1), "DC-DC\n24V->19V/12V/5V", LORANGE, ORANGE, 12)
box(s, Inches(5.6), Inches(1.5), Inches(2.6), Inches(1.0), "Jetson Orin Nano\n(19V)", LGREEN, GREEN, 12)
box(s, Inches(5.6), Inches(3.0), Inches(2.6), Inches(1.0), "STM32 MCU\n(5V)", LBLUE, BLUE, 12)
box(s, Inches(5.6), Inches(4.5), Inches(2.6), Inches(1.0), "Motor driver\n(12-24V) -> 2x DC motor", LORANGE, ORANGE, 11)
box(s, Inches(9.0), Inches(1.2), Inches(3.4), Inches(0.85), "RealSense D435i (USB)", LGREEN, GREEN, 12)
box(s, Inches(9.0), Inches(2.2), Inches(3.4), Inches(0.85), "RPLiDAR A2 (USB)", LGREEN, GREEN, 12)
box(s, Inches(9.0), Inches(3.15), Inches(3.4), Inches(0.85), "SCD30 / BNO055 (I2C)", LBLUE, BLUE, 12)
box(s, Inches(9.0), Inches(4.1), Inches(3.4), Inches(0.85), "Encoders + E-stop (GPIO)", LBLUE, BLUE, 11)
box(s, Inches(9.0), Inches(5.05), Inches(3.4), Inches(0.85), "LED light + pan-tilt (PWM)", LBLUE, BLUE, 11)
arrow(s, Inches(2.7), Inches(3.35), Inches(0.35), Inches(0.45))
arrow(s, Inches(5.25), Inches(3.35), Inches(0.4), Inches(0.45))
caption(s, "Power rails (left) feed compute (centre); sensors/actuators (right) connect via USB / I2C / GPIO / PWM")

# 12. Stage 4 - software architecture
s = new("Software Architecture", "Stage 4")
box(s, Inches(0.6), Inches(1.7), Inches(2.8), Inches(3.6),
    "PERCEPTION\n\n• YOLO26 detect+count\n• SAM size measure\n• VLM anomaly grade\n• RGB-D depth", LGREEN, GREEN, 13, NAVY)
box(s, Inches(3.7), Inches(1.7), Inches(2.8), Inches(3.6),
    "NAVIGATION\n\n• LiDAR SLAM\n• Line following\n• Path planning\n• Obstacle stop", LBLUE, BLUE, 13, NAVY)
box(s, Inches(6.8), Inches(1.7), Inches(2.8), Inches(3.6),
    "MISSION CTRL\n\n• Patrol scheduler\n• Rack waypoints\n• Charge / dock\n• Safety FSM", LORANGE, ORANGE, 13, NAVY)
box(s, Inches(9.9), Inches(1.7), Inches(2.8), Inches(3.6),
    "AIoT / CLOUD\n\n• MQTT publish\n• Dashboard\n• Alerts\n• Yield database", RGBColor(0xFF,0xF2,0xCC), ORANGE, 13, NAVY)
for xx in (3.4, 6.5, 9.6):
    arrow(s, Inches(xx), Inches(3.2), Inches(0.35), Inches(0.5))
caption(s, "Reuses Exercise-9 scripts: detect.py (tasks 1&2), measure_size.py (task 3), agentic_detect.py (VLM)")

# 13. System architecture (capstone)
s = new("System Architecture (AIoT End-to-End)", "Capstone")
box(s, Inches(0.5), Inches(2.6), Inches(3.0), Inches(1.8),
    "EDGE — Robot\n\nSensors + Jetson\nYOLO26 / SAM / VLM\nNavigation", LGREEN, GREEN, 13, NAVY)
box(s, Inches(4.6), Inches(2.6), Inches(3.0), Inches(1.8),
    "FOG — Gateway\n\nWi-Fi AP / MQTT broker\nBuffer + edge rules", LBLUE, BLUE, 13, NAVY)
box(s, Inches(8.7), Inches(2.6), Inches(3.0), Inches(1.8),
    "CLOUD\n\nDatabase + dashboard\nYield analytics\nAlerts to phone", LORANGE, ORANGE, 13, NAVY)
arrow(s, Inches(3.6), Inches(3.25), Inches(0.9), Inches(0.5))
arrow(s, Inches(7.7), Inches(3.25), Inches(0.9), Inches(0.5))
box(s, Inches(4.6), Inches(5.0), Inches(3.0), Inches(0.9), "Farmer / Manager\n(web + mobile app)", LGRAY, GRAY, 12, NAVY)
arrow(s, Inches(9.9), Inches(4.4), Inches(0.4), Inches(0.6), down=True)
bullets(s, [("Three-tier AIoT: low-latency AI at the edge, reliable transport in the fog, analytics & UX in the cloud.", 0, NAVY, True)],
        t=Inches(1.55), h=Inches(0.9))

# 14. Bill of materials (capstone)
s = new("Bill of Materials (BOM)", "Capstone")
table(s, ["Item", "Qty", "Unit (USD)", "Subtotal"], [
    ["Jetson Orin Nano 8GB dev kit", "1", "500", "500"],
    ["Intel RealSense D435i", "1", "300", "300"],
    ["RPLiDAR A2M12", "1", "300", "300"],
    ["Differential-drive base + motors/encoders", "1", "400", "400"],
    ["STM32 MCU + motor driver", "1", "60", "60"],
    ["SCD30 climate sensor + BNO055 IMU", "1", "150", "150"],
    ["24V 20Ah LiFePO4 + BMS + charger/dock", "1", "330", "330"],
    ["DC-DC converters + wiring + e-stop", "1", "120", "120"],
    ["Pan-tilt + LED ring light", "1", "150", "150"],
    ["Frame, wheels, enclosure (IP54), misc", "1", "190", "190"],
    ["TOTAL (hardware)", "", "", "2,500"],
], rh=0.42, fsz=12.5, widths=[6.4, 1.2, 2.0, 2.3], highlight_last=True)
caption(s, "Indicative prices; total hardware BOM ~ USD 2,500 (within the < 3,000 target)")

# 15. Economic evaluation (capstone)
s = new("Economic Evaluation", "Capstone")
table(s, ["Item", "Value (USD/yr)", "Note"], [
    ["Capital cost (one-time)", "2,500", "hardware BOM"],
    ["Software / setup (one-time)", "1,000", "integration + commissioning"],
    ["Operating cost", "-500", "electricity + maintenance + cloud"],
    ["Labor saving", "+11,700", "2 worker-hours/day @ USD 8 saved"],
    ["Yield gain (early anomaly catch)", "+4,000", "~3% loss avoided on output"],
    ["Net annual benefit", "+15,200", "savings - operating cost"],
], rh=0.5, fsz=13, widths=[4.6, 3.0, 4.3], highlight_last=True)
bullets(s, [
    ("Total investment = 3,500 USD  ->  Payback period ~ 2.8 months.", 0, GREEN, True),
    ("5-year net benefit ~ 72,500 USD; ROI strongly positive.", 0, NAVY, True),
], t=Inches(5.3), h=Inches(1.5))

# 16. Final operation (capstone)
s = new("Final Operation", "Capstone")
steps = ["Undock &\nself-test", "Patrol aisles\n(SLAM/line)", "At each rack:\nimage + climate",
         "Edge AI:\ncount/size/grade", "Push data\n& alerts", "Return &\nauto-charge"]
x = Inches(0.45)
for i, st in enumerate(steps):
    box(s, x, Inches(1.7), Inches(1.95), Inches(1.3), st, LBLUE if i % 2 else LGREEN,
        BLUE if i % 2 else GREEN, 11)
    x += Inches(2.05)
    if i < 5:
        arrow(s, x - Inches(0.13), Inches(2.15), Inches(0.28), Inches(0.4))
bullets(s, [
    ("Schedule: 3 automated patrols/day; continuous when racks are near harvest.", 0, NAVY, True),
    ("KPIs: aisle coverage >= 98%, count error <= 10%, size error <= 2 mm, uptime >= 95%.", 0, GRAY, False),
    ("Alerts: contamination / abnormal growth / climate out-of-range -> instant push to manager app.", 0, GRAY, False),
    ("Outputs: per-bag count & size history, yield forecast, harvest-ready bag list.", 0, GREEN, True),
], t=Inches(3.5), h=Inches(3))

# 17. Detection demo (link to Ex.9)
if detect_img:
    s = new("Vision Payload Demo (from Exercise 9)", "Capstone")
    pic = s.shapes.add_picture(detect_img, 0, Inches(1.5), width=Inches(8.2))
    pic.left = int((SW - pic.width) / 2); pic.top = Inches(1.55)
    caption(s, "The robot's onboard YOLO26 model detecting & counting king oyster mushrooms (52 in this rack image)")

# 18. Conclusion
s = new("Conclusion")
bullets(s, [
    ("Completed the full 4-stage engineering design for an AI inspection robot:", 0, NAVY, True),
    ("Stage 1 — requirements list from real grow-house pain points.", 1),
    ("Stage 2 — function structure, morphological matrix, weighted concept selection.", 1),
    ("Stage 3 — modular embodiment & subsystem design.", 1),
    ("Stage 4 — component specs, wiring, software architecture.", 1),
    ("Capstone — system architecture, BOM (~USD 2,500), economics (payback ~2.8 mo), operation.", 0, GREEN, True),
    ("Tightly integrates the Exercise-9 YOLO26 + SAM + LLM/VLM vision system as the robot's brain.", 0, NAVY, True),
], t=Inches(1.6), h=Inches(5.3))

prs.save(str(OUT_PPTX))
print("Saved:", OUT_PPTX, "(", len(prs.slides.__iter__.__self__._sldIdLst), "slides )")
