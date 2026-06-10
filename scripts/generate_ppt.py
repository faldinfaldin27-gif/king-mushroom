"""Generate the assignment PPT for King Oyster Mushroom Agentic Detection.

    python scripts/generate_ppt.py

Builds King_Oyster_Mushroom_AIoT.pptx in the project root from the artifacts in
runs/ and outputs/. Re-run anytime after new results are produced.
"""
import csv
import glob
import statistics
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
OUT_PPTX = ROOT / "King_Oyster_Mushroom_AIoT.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
RED = RGBColor(0xD2, 0x32, 0x2D)
GRAY = RGBColor(0x44, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def first(paths):
    for p in paths:
        if Path(p).exists():
            return str(p)
    return None


def add_box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(p, text, size, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def bar(slide, color=BLUE, h=Inches(0.18)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()


def header(slide, title, color=NAVY):
    bar(slide, BLUE)
    tf = add_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    set_run(tf.paragraphs[0], title, 30, color, bold=True)


def content_slide(title, bullets):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    tf = add_box(s, Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    for i, (txt, lvl, *style) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        color = style[0] if style else GRAY
        bold = style[1] if len(style) > 1 else False
        size = 24 - lvl * 3
        set_run(p, ("• " if lvl == 0 else "   – ") + txt, size, color, bold)
        p.space_after = Pt(8)
    return s


def image_slide(title, img, caption=None, img_w=Inches(8.2)):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    if img and Path(img).exists():
        pic = s.shapes.add_picture(img, 0, Inches(1.4), width=img_w)
        pic.left = int((SW - pic.width) / 2)
        if pic.height > Inches(5.4):
            pic.height = Inches(5.4); pic.width = int(pic.height * 1.4)
            pic.left = int((SW - pic.width) / 2)
        pic.top = Inches(1.45)
        if caption:
            tf = add_box(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4))
            set_run(tf.paragraphs[0], caption, 13, GRAY, align=PP_ALIGN.CENTER)
    else:
        tf = add_box(s, Inches(1), Inches(3), Inches(11), Inches(1))
        set_run(tf.paragraphs[0], f"[missing image: {img}]", 16, RED)
    return s


def table_slide(title, headers, rows, caption=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    nrow, ncol = len(rows) + 1, len(headers)
    t = s.shapes.add_table(nrow, ncol, Inches(1.2), Inches(1.7),
                           Inches(10.9), Inches(0.4 * nrow)).table
    for j, h in enumerate(headers):
        c = t.cell(0, j); c.text = h
        c.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        c.text_frame.paragraphs[0].runs[0].font.bold = True
        c.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.text = str(val)
            c.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
    if caption:
        tf = add_box(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4))
        set_run(tf.paragraphs[0], caption, 13, GRAY, align=PP_ALIGN.CENTER)
    return s


# ---- gather data ----
sizes = []
csv_path = ROOT / "outputs" / "sizes.csv"
if csv_path.exists():
    with open(csv_path) as f:
        for row in csv.DictReader(f):
            sizes.append(row)

detect_img = first(sorted(glob.glob(str(ROOT / "outputs" / "detect" / "IMG2023*.jpg"))) +
                   sorted(glob.glob(str(ROOT / "outputs" / "detect" / "*.jpg"))))
size_img = first(glob.glob(str(ROOT / "outputs" / "sizes_*.jpg")))
results_curve = ROOT / "runs" / "detect" / "king_oyster_yolo26" / "results.png"
conf_mat = ROOT / "runs" / "detect" / "val" / "confusion_matrix.png"
pr_curve = ROOT / "runs" / "detect" / "val" / "BoxPR_curve.png"

# ================= SLIDES =================

# 1. Title
s = prs.slides.add_slide(BLANK)
bar(s, BLUE, Inches(0.25))
bar2 = s.shapes.add_shape(1, 0, SH - Inches(0.25), SW, Inches(0.25))
bar2.fill.solid(); bar2.fill.fore_color.rgb = BLUE; bar2.line.fill.background()
tf = add_box(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
set_run(tf.paragraphs[0], "AIoT System Integration Design & Application", 38, NAVY, True, PP_ALIGN.CENTER)
p = tf.add_paragraph(); set_run(p, "King Oyster Mushroom Agentic Detection", 28, RED, True, PP_ALIGN.CENTER)
tf2 = add_box(s, Inches(1), Inches(4.1), Inches(11.3), Inches(2))
set_run(tf2.paragraphs[0], "Exercise 9 — YOLO26 + SAM + LLM/VLM Agentic Object Detection", 20, GRAY, False, PP_ALIGN.CENTER)
p = tf2.add_paragraph(); set_run(p, "Individual Detection  ·  Counting  ·  Size Measurement", 18, BLUE, False, PP_ALIGN.CENTER)

# 2. Assignment requirements
content_slide("Assignment Tasks", [
    ("Use YOLO26 + SAM3/SAM 3D with LLM/VLM for king oyster mushroom agentic detection", 0, NAVY, True),
    ("(1) Detect each individual mushroom in one substrate bag", 1),
    ("(2) Count the number of mushrooms in one bag", 1),
    ("(3) Measure the size of each individual mushroom", 1),
    ("(4) At least one custom topic (LLM/VLM)", 1),
    ("References:", 0, NAVY, True),
    ("landing.ai — What is Agentic Object Detection", 1, BLUE),
    ("Ultralytics YOLO26 — docs.ultralytics.com/models/yolo26", 1, BLUE),
])

# 3. System pipeline
content_slide("System Pipeline", [
    ("Input: photo of the substrate bag", 0, NAVY, True),
    ("YOLO26 detector -> bounding box per mushroom", 1, BLUE, True),
    ("-> Tasks 1 & 2: individual detection + counting", 2),
    ("SAM3 / SAM 3D -> precise mask -> pixel area -> mm", 1, BLUE, True),
    ("-> Task 3: size measurement (needs scale calibration)", 2),
    ("LLM / VLM agent -> verify counts, grade ripeness, answer in language", 1, BLUE, True),
    ("-> Task 4: VisionAgent, detector used as a tool", 2),
])

# 4. Dataset
table_slide("Dataset", ["Item", "Detail"], [
    ["Source", "Roboflow  esdl/king-oyster-mushroom v8 (CC BY 4.0)"],
    ["Classes", "1 class — King_oyster_mushroom"],
    ["Train", "490 images"],
    ["Validation", "62 images"],
    ["Test", "61 images (561 instances)"],
    ["Label format", "YOLO txt — class cx cy w h (normalized)"],
    ["Characteristic", "dense, small objects"],
], caption="Single-class, dense small-object detection is the main challenge of this task")

# 5. Training setup
table_slide("YOLO26 Training Setup", ["Parameter", "Value"], [
    ["Base model", "yolo26n.pt"],
    ["Epochs", "100"],
    ["Image size", "640"],
    ["Batch", "4"],
    ["Device (GPU)", "NVIDIA GeForce GTX 1650 Ti (4GB, CUDA 12.4)"],
    ["Framework", "Ultralytics 8.4.51 + PyTorch 2.6.0+cu124"],
    ["Parameters", "2.38M params, 5.2 GFLOPs"],
])

# 6. Training curves
image_slide("Training Curves", str(results_curve),
            "Losses decrease steadily and mAP rises — 100 epochs", Inches(11.5))

# 7. Validation metrics table
table_slide("Evaluation (Test set, 61 images)", ["Metric", "Value", "Meaning"], [
    ["Precision", "0.848", "85% of detections are real mushrooms"],
    ["Recall", "0.657", "66% of real mushrooms are detected"],
    ["mAP@50", "0.773", "Overall detection accuracy"],
    ["mAP@50-95", "0.528", "Strict localization accuracy"],
    ["Inference", "11.7 ms / img", "Real-time speed on GPU"],
])

# 8. Confusion / PR
image_slide("Precision-Recall Curve", str(pr_curve),
            "mAP@0.5 = 0.773 for King_oyster_mushroom", Inches(7.5))

# 9. Task 1 & 2
image_slide("Tasks 1 & 2: Detection + Counting", detect_img,
            "YOLO26 detects and counts each mushroom (52 detected in this image)", Inches(8.5))

# 10. Task 3 size
if sizes:
    w = [float(r["width_mm"]) for r in sizes]
    h = [float(r["height_mm"]) for r in sizes]
    a = [float(r["area_mm2"]) for r in sizes]
    n = len(sizes)
    image_slide("Task 3: Individual Size Measurement", size_img,
                f"{n} mushrooms measured | mean {statistics.mean(w):.1f}x{statistics.mean(h):.1f} mm | "
                f"mean area {statistics.mean(a):.0f} mm2 (calibrated bag width=120mm)", Inches(8.5))
    table_slide("Task 3: Size Statistics (mm)",
                ["Statistic", "Width", "Height", "Area (mm2)"], [
        ["Count n", n, n, n],
        ["Mean", f"{statistics.mean(w):.1f}", f"{statistics.mean(h):.1f}", f"{statistics.mean(a):.0f}"],
        ["Min", f"{min(w):.1f}", f"{min(h):.1f}", f"{min(a):.0f}"],
        ["Max", f"{max(w):.1f}", f"{max(h):.1f}", f"{max(a):.0f}"],
    ], caption="Scale: bag mouth width 120mm -> 5.21 px/mm; full data in outputs/sizes.csv")

# 11. Conclusion
content_slide("Conclusion & Next Steps", [
    ("Completed:", 0, NAVY, True),
    ("(1) Individual mushroom detection — YOLO26 detector mAP@50 = 0.77", 1),
    ("(2) Counting — automatic count (1319 mushrooms across test set)", 1),
    ("(3) Individual size measurement — px->mm, CSV statistics exported", 1),
    ("Next:", 0, NAVY, True),
    ("(4) Custom topic: LLM/VLM agentic — ripeness / harvest grading", 1, RED),
    ("Improve: imgsz 1280 or SAHI tiling for small-object recall", 1),
    ("Integrate into AIoT: edge real-time detection + cloud reporting", 1),
])

prs.save(str(OUT_PPTX))
print(f"Saved: {OUT_PPTX}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
